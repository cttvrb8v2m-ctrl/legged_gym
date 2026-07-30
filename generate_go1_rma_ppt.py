#!/usr/bin/env python3
"""Generate a plain 10-slide weekly group-meeting report for the Go1 project."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "Go1_组会汇报_本周进展_10页.pptx"

W = 13.333
H = 7.5
FONT = "Microsoft YaHei"
NAVY = "17365D"
BLUE = "2F5597"
TEXT = "222222"
GRAY = "666666"
LINE = "B7C3D0"
LIGHT = "F3F5F7"
PALE_BLUE = "EAF0F7"
WHITE = "FFFFFF"
RED = "9C2F2F"
GREEN = "3F6B4F"


def rgb(value):
    return RGBColor.from_string(value)


def set_run(run, size=20, color=TEXT, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=20,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.04,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    set_run(run, size, color, bold)
    return box


def add_rich_line(slide, parts, x, y, w, h, size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    paragraph = frame.paragraphs[0]
    for text, bold, color in parts:
        run = paragraph.add_run()
        run.text = text
        set_run(run, size, color, bold)
    return box


def add_bullets(slide, items, x, y, w, h, size=20, gap=9):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(gap)
        paragraph.line_spacing = 1.12
        marker = paragraph.add_run()
        marker.text = "•  "
        set_run(marker, size, NAVY, True)
        body = paragraph.add_run()
        body.text = item
        set_run(body, size, TEXT, False)
    return box


def add_rect(slide, x, y, w, h, fill=WHITE, line=LINE, width=1.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(width)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, width=1.5):
    arrow = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    arrow.line.color.rgb = rgb(color)
    arrow.line.width = Pt(width)
    arrow.line.end_arrowhead = True
    return arrow


def add_header(slide, title, page):
    add_text(slide, title, 0.72, 0.42, 11.7, 0.52, 30, NAVY, True)
    line = add_rect(slide, 0.72, 1.08, 11.9, 0.018, NAVY, NAVY, 0)
    line.line.fill.background()
    add_text(slide, str(page), 12.0, 0.48, 0.5, 0.30, 11, GRAY, False,
             PP_ALIGN.RIGHT)


def add_footer(slide, text="Unitree Go1 阶段工作汇报"):
    add_text(slide, text, 0.74, 7.12, 5.6, 0.18, 9, GRAY)


def add_table(slide, data, x, y, w, h, widths=None, font_size=17):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(
        rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)
    ).table
    if widths:
        for index, col_width in enumerate(widths):
            table.columns[index].width = Inches(col_width)
    for row_index, row in enumerate(data):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(PALE_BLUE if row_index == 0 else WHITE)
            cell.border_left = None
            frame = cell.text_frame
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            paragraph = frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                set_run(
                    run,
                    font_size,
                    NAVY if row_index == 0 else TEXT,
                    row_index == 0,
                )
    return table


def add_flow_box(slide, text, x, y, w, h, fill=WHITE, size=19):
    add_rect(slide, x, y, w, h, fill, NAVY, 1.2)
    add_text(
        slide, text, x + 0.08, y + 0.08, w - 0.16, h - 0.16,
        size, NAVY, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE,
    )


def new_slide(prs, title, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)
    add_header(slide, title, page)
    add_footer(slide)
    return slide


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # 1. Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)
    add_rect(slide, 0.0, 0.0, 0.18, H, NAVY, NAVY, 0)
    add_text(slide, "Unitree Go1 阶段工作汇报", 1.05, 1.55, 10.8, 0.72,
             31, NAVY, True)
    add_text(slide, "连续台阶能力与高速后腿足宽调整", 1.05, 2.50, 9.6, 0.56,
             23, TEXT)
    add_rect(slide, 1.05, 3.43, 10.75, 0.015, LINE, LINE, 0)
    add_text(slide, "本周组会汇报", 1.05, 3.86, 3.5, 0.38, 20, GRAY)
    add_text(slide, "2026.07", 1.05, 5.85, 2.0, 0.30, 17, GRAY)
    add_text(slide, "强化学习 / 课程学习 / RMA结构检查", 7.55, 5.82, 4.25,
             0.34, 15, GRAY, False, PP_ALIGN.RIGHT)

    # 2. Weekly overview
    slide = new_slide(prs, "本周工作概述", 2)
    add_bullets(
        slide,
        [
            "本周我继续以 model_808 作为稳定基线，没有再沿用成功率下降的 model_810 分支。",
            "我重点处理了高速运动时后腿逐渐并拢的问题，并补充了足宽、扭矩和摔倒率评估。",
            "台阶部分主要重新整理了 0.150、0.155 和 0.160 m 的失败类型，确认后腿越过立面仍是主要瓶颈。",
            "我也检查了当前 RMA 的实现和加载流程，但本周实验还没有验证出明确、稳定的 RMA 增益。",
        ],
        0.92, 1.46, 11.45, 3.75, 20,
    )
    add_rect(slide, 0.92, 5.48, 11.42, 0.88, LIGHT, LINE, 1)
    add_rich_line(
        slide,
        [
            ("本周实际有效的改进：", True, NAVY),
            ("速度门控的后髋残差；", False, TEXT),
            ("台阶已有能力主要来自课程学习和台阶事件奖励。", False, TEXT),
        ],
        1.15, 5.72, 10.95, 0.40, 19,
    )

    # 3. High-speed problem
    slide = new_slide(prs, "高速足宽：我先确认了问题出现在哪里", 3)
    add_bullets(
        slide,
        [
            "速度升高以后，前腿平均足宽基本不变，明显收窄的是后腿支撑宽度。",
            "rear_slip 和 torque_saturation 基本为 0，所以这不是打滑或力矩上限直接造成的。",
            "从关节动作看，原 Actor 在高速步态下没有给后髋足够的向外落脚量。",
        ],
        0.86, 1.35, 11.55, 2.25, 19,
    )
    add_text(slide, "model_808 后腿接触宽度 P05", 0.92, 3.72, 4.2, 0.34,
             19, NAVY, True)
    add_table(
        slide,
        [
            ["速度指令", "2.5 m/s", "3.0 m/s", "3.5 m/s"],
            ["后腿宽度 P05", "10.04 cm", "8.71 cm", "7.42 cm"],
        ],
        0.92, 4.18, 8.15, 1.15, [2.15, 2.0, 2.0, 2.0], 18,
    )
    add_rect(slide, 9.45, 4.18, 2.85, 1.15, LIGHT, LINE, 1)
    add_text(slide, "速度越高，\n后腿窄尾部越明显", 9.68, 4.39, 2.40, 0.70,
             19, RED, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

    # 4. High-speed method
    slide = new_slide(prs, "高速足宽：采用速度门控的后髋残差", 4)
    add_bullets(
        slide,
        [
            "我保留原 Actor，不重新训练整套步态，只在左右后髋动作上叠加一个较小的向外残差。",
            "残差在速度指令低于 2.0 m/s 时为 0，到高速区间后再逐渐生效。",
            "训练时冻结 Actor、RMA 和动作标准差，只更新后髋残差模块与 Critic。",
            "选择共享的左右对称输出，主要是为了避免两条后腿在高速下产生新的横向偏置。",
        ],
        0.86, 1.34, 11.55, 2.65, 19,
    )
    add_flow_box(slide, "当前观测\n与速度指令", 1.05, 4.45, 2.25, 0.95, LIGHT)
    add_arrow(slide, 3.30, 4.93, 4.10, 4.93)
    add_flow_box(slide, "后髋残差模块", 4.10, 4.45, 2.35, 0.95, PALE_BLUE)
    add_arrow(slide, 6.45, 4.93, 7.25, 4.93)
    add_flow_box(slide, "速度门控", 7.25, 4.45, 1.95, 0.95, WHITE)
    add_arrow(slide, 9.20, 4.93, 10.00, 4.93)
    add_flow_box(slide, "左右后髋\n向外修正", 10.00, 4.45, 2.20, 0.95, PALE_BLUE)
    add_text(slide, "楼梯速度下残差严格为 0", 4.15, 5.78, 5.0, 0.34,
             18, GREEN, True, PP_ALIGN.CENTER)

    # 5. Mechanism and width result
    slide = new_slide(prs, "高速足宽：作用机理和当前结果", 5)
    add_flow_box(slide, "原始 Actor 动作", 0.92, 1.48, 2.55, 0.90, WHITE)
    add_text(slide, "+", 3.62, 1.72, 0.45, 0.34, 24, NAVY, True,
             PP_ALIGN.CENTER)
    add_flow_box(slide, "高速后髋残差", 4.16, 1.48, 2.60, 0.90, PALE_BLUE)
    add_arrow(slide, 6.95, 1.93, 8.08, 1.93)
    add_flow_box(slide, "后腿向外落脚", 8.20, 1.48, 2.65, 0.90, WHITE)
    add_text(slide, "支撑宽度增加", 11.10, 1.70, 1.45, 0.40, 18, GREEN, True,
             PP_ALIGN.CENTER)
    add_table(
        slide,
        [
            ["速度", "model_808 P05", "model_880 P05", "变化"],
            ["2.5 m/s", "10.04 cm", "11.89 cm", "+1.85 cm"],
            ["3.0 m/s", "8.71 cm", "10.67 cm", "+1.96 cm"],
            ["3.5 m/s", "7.42 cm", "9.31 cm", "+1.89 cm"],
        ],
        0.92, 3.03, 8.75, 2.45, [1.65, 2.35, 2.35, 2.40], 17,
    )
    add_bullets(
        slide,
        [
            "前腿平均足宽基本不变。",
            "3.5 m/s 三个种子共 300 次测试没有摔倒。",
            "当前仍有部分着地点过窄或过宽，还不能认为足宽已经完全稳定。",
        ],
        9.95, 3.05, 2.38, 2.75, 16, 6,
    )

    # 6. Stair problem
    slide = new_slide(prs, "台阶能力：主要问题仍是后腿越不过立面", 6)
    add_bullets(
        slide,
        [
            "前腿通常能够先上阶，但后腿起摆后不能稳定越过立面并落到踏面。",
            "随着台阶从 0.150 m 增加到 0.160 m，rear_not_clear 比例明显增加。",
            "后腿足宽与失败没有正相关，因此高速外展方案不适合直接搬到楼梯上。",
        ],
        0.86, 1.35, 11.55, 2.30, 19,
    )
    add_table(
        slide,
        [
            ["台阶高度", "成功率", "平均通过级数", "rear failure"],
            ["0.150 m", "37.0%", "6.37", "61.7%"],
            ["0.155 m", "17.0%", "4.88", "82.0%"],
            ["0.160 m", "10.3%", "3.57", "89.0%"],
        ],
        0.92, 3.98, 8.45, 2.28, [1.8, 1.8, 2.4, 2.45], 17,
    )
    add_rect(slide, 9.72, 3.98, 2.60, 2.28, LIGHT, LINE, 1)
    add_text(slide, "我目前的判断", 9.95, 4.25, 2.15, 0.32, 18, NAVY, True,
             PP_ALIGN.CENTER)
    add_text(slide, "瓶颈主要在后腿的\n前向越阶时机、\n垂直净空和踏面落点，\n不是横向足宽。",
             9.98, 4.78, 2.08, 1.15, 17, TEXT, False,
             PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

    # 7. Stair mechanism
    slide = new_slide(prs, "台阶能力：课程学习和事件奖励的作用", 7)
    add_bullets(
        slide,
        [
            "台阶能力不是本周临时加入的动作补丁，而是在逐级增加台阶高度的课程中形成的。",
            "训练中使用台阶进度、登顶、后腿目标进度等事件奖励，帮助策略把动作和台阶阶段对应起来。",
            "当前有效机理是先建立前腿支撑，再完成重心转移，为后腿起摆和落脚留出空间。",
        ],
        0.86, 1.34, 11.55, 2.35, 19,
    )
    labels = ["前腿上阶", "重心转移", "后腿起摆", "后腿落到踏面"]
    xs = [0.92, 4.00, 7.08, 10.16]
    for index, (label, x) in enumerate(zip(labels, xs)):
        add_flow_box(slide, label, x, 4.35, 2.15, 0.95,
                     PALE_BLUE if index in (0, 3) else WHITE, 18)
        if index < 3:
            add_arrow(slide, x + 2.15, 4.83, x + 2.95, 4.83)
    add_text(slide, "课程学习决定“能不能上”，事件奖励帮助动作与台阶阶段对齐。",
             1.22, 5.82, 10.90, 0.40, 19, NAVY, True, PP_ALIGN.CENTER)

    # 8. RMA understanding
    slide = new_slide(prs, "RMA：本周没有验证出明确增益", 8)
    add_bullets(
        slide,
        [
            "我检查了 history encoder、FiLM、alpha、优化器和 checkpoint 加载流程，确认这些模块在代码中能够参与前向计算。",
            "但当前结构没有完整的 privileged teacher latent 和 latent imitation loss，因此还不是真正的两阶段 RMA。",
            "打开 alpha 后动作会发生明显变化，而现有实验没有给出稳定的成功率提升，所以本周继续使用 alpha=0。",
            "因此，本周高速足宽改善不能归因于 RMA；它来自单独训练的速度门控后髋残差。",
        ],
        0.86, 1.34, 11.55, 3.45, 19,
    )
    add_rect(slide, 0.92, 5.15, 11.42, 0.92, LIGHT, LINE, 1)
    add_rich_line(
        slide,
        [
            ("后续如果继续做 RMA：", True, NAVY),
            ("先训练真实台阶高度的 teacher latent，再冻结 teacher，用历史观测拟合 latent。", False, TEXT),
        ],
        1.16, 5.39, 10.95, 0.44, 18,
    )

    # 9. Experiments and takeaways
    slide = new_slide(prs, "本周实验与认识", 9)
    add_bullets(
        slide,
        [
            "model_880 在 3.5 m/s 下把后腿接触宽度 P05 从 7.42 cm 提高到 9.31 cm，三种子共 300 次没有摔倒。",
            "前腿平均足宽约 15.7 cm，前腿没有因为后髋残差发生明显横向变化。",
            "我没有选择最后的 model_958，因为它的残差硬饱和率已经达到 27.6%，稳定余量较小。",
            "楼梯速度下残差严格为 0，model_880 的 Actor、RMA 和动作标准差与 model_808 保持一致。",
            "目前结果说明“小范围、带门控的动作修正”有效，但还需要继续处理着地点分布偏宽和偏窄并存的问题。",
        ],
        0.86, 1.34, 11.55, 4.30, 19,
    )
    add_rect(slide, 0.92, 5.90, 11.42, 0.54, PALE_BLUE, LINE, 1)
    add_text(slide, "本周结论以稳定性和可复现评估为准，不把尚未验证的 RMA 效果写成结果。",
             1.14, 6.03, 10.98, 0.28, 18, NAVY, True, PP_ALIGN.CENTER)

    # 10. Summary and deliverables
    slide = new_slide(prs, "总结与交付", 10)
    add_bullets(
        slide,
        [
            "本周我完成了高速后腿足宽问题的定位、两轮残差训练和多速度评估，当前推荐使用 model_880。",
            "台阶部分保留 model_808 的原策略能力，不让高速足宽修正干扰低速登阶。",
            "RMA 本周只完成了结构检查和风险确认，还没有形成可以单独报告的性能增益。",
            "下一步优先继续分析后腿起摆时机和踏面落点，再单独设计真正的两阶段 RMA 实验。",
        ],
        0.86, 1.34, 11.55, 3.25, 19,
    )
    add_text(slide, "当前交付模型", 0.92, 4.90, 2.2, 0.34, 19, NAVY, True)
    add_rect(slide, 0.92, 5.38, 11.42, 0.78, LIGHT, LINE, 1)
    add_text(
        slide,
        "logs/rough_go1/Jul29_11-39-32_highspeed_rear_hip_contact_band100_from858/model_880.pt",
        1.14, 5.60, 10.95, 0.34, 16, TEXT, False, PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(slide, "说明：演示时 residual_max_action=0.12，RMA alpha=0。",
             1.15, 6.48, 10.95, 0.30, 17, GRAY, False, PP_ALIGN.CENTER)

    prs.core_properties.title = "Unitree Go1 阶段工作汇报"
    prs.core_properties.subject = "连续台阶能力与高速后腿足宽调整"
    prs.core_properties.author = "Go1 项目组会汇报"
    prs.core_properties.keywords = "Go1, PPO, RMA, 课程学习, 高速足宽"
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    result = create_presentation()
    print(f"Generated: {result}")
